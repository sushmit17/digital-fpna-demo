"""
Build the 5-slide S&OP PowerPoint deck from consolidated P&L data.
Style: Clean corporate white matching the design document.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x38, 0x64)
BLUE   = RGBColor(0x2E, 0x5F, 0xA3)
TEAL   = RGBColor(0x00, 0xB8, 0xA0)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xF2, 0xF4, 0xF8)
MGREY  = RGBColor(0xD6, 0xDC, 0xE4)
DGREY  = RGBColor(0x6B, 0x7A, 0x8D)
GREEN  = RGBColor(0x1A, 0x6B, 0x3C)
RED    = RGBColor(0xC0, 0x39, 0x2B)
DARK   = RGBColor(0x1A, 0x1F, 0x3E)

SLD_W = Inches(13.33)
SLD_H = Inches(7.5)

def rgb_hex(r):
    return f"{r[0]:02X}{r[1]:02X}{r[2]:02X}"

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w or 1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h,
                font_size=12, bold=False, color=DARK,
                align=PP_ALIGN.LEFT, font_name="Calibri", wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def slide_header(slide, title, subtitle=None):
    """Navy top bar with title."""
    add_rect(slide, 0, 0, 13.33, 1.05, NAVY)
    add_rect(slide, 0, 1.05, 13.33, 0.04, TEAL)
    add_textbox(slide, title, 0.45, 0.12, 10, 0.65,
                font_size=26, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle, 0.45, 0.7, 10, 0.32,
                    font_size=12, color=RGBColor(0xC5, 0xD8, 0xEC))
    # Page footer
    add_rect(slide, 0, 7.2, 13.33, 0.3, RGBColor(0xEE, 0xF2, 0xF8))
    add_textbox(slide, "Digital FP&A Manager  |  Commercial Affiliate  |  LBE FY2026  |  Confidential",
                0.3, 7.22, 10, 0.25, font_size=8, color=DGREY)
    add_textbox(slide, "EUR millions", 11.8, 7.22, 1.3, 0.25,
                font_size=8, color=DGREY, align=PP_ALIGN.RIGHT)


def fmt(v, show_sign=False):
    """Format a number as EUR millions string."""
    if v is None:
        return "-"
    if show_sign and v > 0:
        return f"+{v:,.1f}"
    if v == 0:
        return "-"
    return f"({abs(v):,.1f})" if v < 0 else f"{v:,.1f}"

def var_color(v, is_cost=False):
    """Return GREEN or RED based on favourable/unfavourable."""
    if v == 0 or v is None:
        return DGREY
    fav = (v > 0 and not is_cost) or (v < 0 and is_cost)
    return GREEN if fav else RED


# ── P&L table helper ──────────────────────────────────────────────────────────
PNL_ROWS = [
    ("Net Sales",                       False, "total"),
    ("Distribution Margin",             False, "total"),
    ("Total R&D",                       True,  "total"),
    ("Marketing",                       True,  "detail"),
    ("Other Advertising and Promotion", True,  "detail"),
    ("Sales Force",                     True,  "detail"),
    ("General Admin",                   True,  "detail"),
    ("Total SG&A",                      True,  "total"),
    ("Division Margin",                 False, "total"),
]

def draw_pnl_table(slide, data, cols_def, top=1.2, left=0.3,
                   col_widths=None, title=None):
    """
    cols_def: list of (header, data_key, show_sign, is_cost_col)
    """
    n_cols = len(cols_def) + 1  # +1 for label col
    if col_widths is None:
        label_w = 3.0
        data_w  = (13.33 - left * 2 - label_w) / len(cols_def)
        col_widths = [label_w] + [data_w] * len(cols_def)

    row_h = 0.38
    hdr_h = 0.45

    if title:
        add_textbox(slide, title, left, top - 0.38, 12, 0.35,
                    font_size=11, bold=True, color=NAVY)

    # Header row
    x = left
    hdrs = ["P&L Line"] + [c[0] for c in cols_def]
    for i, (hdr, cw) in enumerate(zip(hdrs, col_widths)):
        fill_c = NAVY if i == 0 else (TEAL if "Var" in hdr else BLUE)
        add_rect(slide, x, top, cw - 0.02, hdr_h, fill_c)
        add_textbox(slide, hdr, x + 0.04, top + 0.04, cw - 0.08, hdr_h - 0.08,
                    font_size=9, bold=True, color=WHITE,
                    align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
        x += cw

    # Data rows
    for ri, (line, is_cost, ltype) in enumerate(PNL_ROWS):
        if line not in data:
            continue
        y = top + hdr_h + ri * row_h
        x = left
        row_fill = RGBColor(0xEE, 0xF2, 0xF8) if ltype == "total" else (
            LGREY if ri % 2 == 0 else WHITE)
        is_bold = ltype == "total"

        for ci, cw in enumerate(col_widths):
            add_rect(slide, x, y, cw - 0.02, row_h - 0.02, row_fill)
            if ci == 0:
                indent = 0.06 if ltype == "total" else 0.2
                add_textbox(slide, line, x + indent, y + 0.05,
                            cw - indent - 0.04, row_h - 0.1,
                            font_size=9, bold=is_bold, color=NAVY if is_bold else DARK)
            else:
                _, data_key, show_sign, is_cost_col = cols_def[ci - 1]
                v = data.get(line, {}).get(data_key, 0)
                txt = fmt(v, show_sign=show_sign)
                clr = var_color(v, is_cost) if show_sign else (
                    NAVY if is_bold else DARK)
                add_textbox(slide, txt, x + 0.02, y + 0.05,
                            cw - 0.06, row_h - 0.1,
                            font_size=9, bold=is_bold, color=clr,
                            align=PP_ALIGN.RIGHT)
            x += cw


# ── SLIDE BUILDERS ────────────────────────────────────────────────────────────

def build_slide1_summary(prs, data, narrative):
    """Slide 1: Executive Summary — KPI callouts + narrative."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide_header(slide, "Executive Summary", "LBE FY2026 — Affiliate P&L Overview")

    # KPI cards (Net Sales, Distribution Margin, Division Margin — LBE vs Budget)
    kpis = [
        ("Net Sales",         "Var vs Budget", False),
        ("Distribution Margin","Var vs Budget", False),
        ("Division Margin",   "Var vs Budget", False),
    ]
    card_w = 3.8
    for i, (line, var_key, is_cost) in enumerate(kpis):
        lbe_v = data.get(line, {}).get("LBE FY2026", 0)
        var_v = data.get(line, {}).get(var_key, 0)
        xl = 0.3 + i * (card_w + 0.25)

        add_rect(slide, xl, 1.2, card_w, 1.5,
                 RGBColor(0xEE, 0xF2, 0xF8))
        add_rect(slide, xl, 1.2, card_w, 0.06, TEAL)
        add_textbox(slide, line, xl + 0.15, 1.28, card_w - 0.3, 0.35,
                    font_size=10, bold=True, color=NAVY)
        add_textbox(slide, f"{fmt(lbe_v)}", xl + 0.15, 1.6, card_w - 0.3, 0.55,
                    font_size=28, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
        clr = var_color(var_v, is_cost)
        sign = "+" if var_v > 0 else ""
        add_textbox(slide, f"{sign}{var_v:,.1f} vs Budget",
                    xl + 0.15, 2.1, card_w - 0.3, 0.35,
                    font_size=10, color=clr)

    # Narrative text box
    add_rect(slide, 0.3, 2.85, 12.73, 0.04, TEAL)
    add_textbox(slide, "VARIANCE NARRATIVE", 0.3, 2.95, 5, 0.3,
                font_size=9, bold=True, color=NAVY)
    add_textbox(slide, narrative, 0.3, 3.25, 12.7, 3.7,
                font_size=10.5, color=DARK, wrap=True)



def build_slide2_volumes(prs, unit_data: list):
    """
    Slide 2 (new): Product Volume by Therapeutic Area.
    Shows LBE vs Budget vs Prior Year volumes (000 units) for every product,
    grouped by TA, in a clean grouped table layout.
    unit_data: list of parsed unit dicts (with 'products' key for Sales units)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Product Volume by Therapeutic Area",
                 "LBE FY2026 vs Budget & Prior Year  —  Volumes in 000 units  (shown as 000s)")

    # Filter to Sales units that have products
    sales_units = [u for u in unit_data if u.get("products")]
    if not sales_units:
        add_textbox(slide, "No product volume data available.",
                    0.5, 2.0, 12, 1.0, font_size=14, color=DGREY)
        return

    # ── Layout: one column per TA, table inside each ──────────────────────────
    n_units = len(sales_units)
    col_w   = (13.33 - 0.4) / n_units        # width of each TA block
    left    = 0.2
    top     = 1.25
    hdr_h   = 0.38
    row_h   = 0.36
    ta_colors = [BLUE, TEAL, RGBColor(0x8B, 0x5C, 0xF6)]

    for ui, unit in enumerate(sales_units):
        xl = left + ui * col_w
        block_w = col_w - 0.15
        ta_color = ta_colors[ui % len(ta_colors)]

        # TA header band
        add_rect(slide, xl, top, block_w, hdr_h + 0.04, ta_color)
        add_textbox(slide, unit["unit"].upper(),
                    xl + 0.08, top + 0.05, block_w - 0.16, hdr_h - 0.1,
                    font_size=11, bold=True, color=WHITE)

        # Sub-header row: Product | LBE | Budget | PY | vs Bgt
        sub_top = top + hdr_h + 0.04
        sub_cols_w = [block_w * f for f in [0.35, 0.14, 0.14, 0.14, 0.23]]
        sub_hdrs   = ["Product", "LBE", "Bgt", "PY", "vs Bgt"]
        x = xl
        for sci, (sh, sw) in enumerate(zip(sub_hdrs, sub_cols_w)):
            add_rect(slide, x, sub_top, sw - 0.02, 0.3,
                     RGBColor(0xEE, 0xF2, 0xF8))
            add_textbox(slide, sh, x + 0.02, sub_top + 0.03,
                        sw - 0.04, 0.24,
                        font_size=8, bold=True, color=NAVY,
                        align=PP_ALIGN.LEFT if sci == 0 else PP_ALIGN.CENTER)
            x += sw

        # Data rows
        for ri, prod in enumerate(unit["products"]):
            ry = sub_top + 0.3 + ri * row_h
            row_fill = RGBColor(0xF8, 0xF9, 0xFB) if ri % 2 == 0 else WHITE
            lbe_v = prod.get("lbe_vol", 0)
            bgt_v = prod.get("bgt_vol", 0)
            py_v  = prod.get("py_vol", 0)
            var_v = round(lbe_v - bgt_v, 0)

            # Truncate product name if too long
            short_name = prod["name"]
            if "(" in short_name:
                short_name = short_name[:short_name.index("(")].strip()

            x = xl
            def fv(n, sign=False):
                # Format as thousands with one decimal: 42800 -> "42.8"
                s = f"{abs(n)/1000:.1f}"
                if sign: return ("+" if n >= 0 else "") + f"{n/1000:.1f}"
                return s
            vals = [short_name, fv(lbe_v), fv(bgt_v), fv(py_v), fv(var_v, sign=True)]
            for ci, (val, sw) in enumerate(zip(vals, sub_cols_w)):
                add_rect(slide, x, ry, sw - 0.02, row_h - 0.02, row_fill)
                is_var = ci == 4
                if is_var:
                    clr = GREEN if var_v >= 0 else RED
                else:
                    clr = NAVY if ci == 0 else DARK
                add_textbox(slide, val,
                            x + 0.02, ry + 0.04, sw - 0.06, row_h - 0.08,
                            font_size=9, bold=(ci == 0), color=clr,
                            align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
                x += sw

        # Total row
        total_row = sub_top + 0.3 + len(unit["products"]) * row_h
        t_lbe = sum(p.get("lbe_vol", 0) for p in unit["products"])
        t_bgt = sum(p.get("bgt_vol", 0) for p in unit["products"])
        t_py  = sum(p.get("py_vol",  0) for p in unit["products"])
        t_var = round(t_lbe - t_bgt, 0)

        x = xl
        def fv2(n, sign=False):
            if sign: return ("+" if n >= 0 else "") + f"{n/1000:.1f}"
            return f"{abs(n)/1000:.1f}"
        tot_vals = ["Total", fv2(t_lbe), fv2(t_bgt), fv2(t_py), fv2(t_var, sign=True)]
        for ci, (val, sw) in enumerate(zip(tot_vals, sub_cols_w)):
            add_rect(slide, x, total_row, sw - 0.02, row_h - 0.02,
                     RGBColor(0xEE, 0xF2, 0xF8))
            is_var = ci == 4
            clr = (GREEN if t_var >= 0 else RED) if is_var else NAVY
            add_textbox(slide, val,
                        x + 0.02, total_row + 0.04, sw - 0.06, row_h - 0.08,
                        font_size=9, bold=True, color=clr,
                        align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
            x += sw


def build_slide2_sales(prs, data):
    """Slide 2: Sales Forecast by Therapeutic Area."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Sales Forecast by Therapeutic Area",
                 "Net Sales & Distribution Margin — LBE FY2026")

    # Sales bar chart (simulated with rectangles — no chart library needed)
    tas = [
        ("Immunology",  data.get("_immunology_sales", 312.4),
                        data.get("_immunology_budget", 295.0)),
        ("Oncology",    data.get("_oncology_sales",   187.6),
                        data.get("_oncology_budget",  210.0)),
        ("Neurology",   data.get("_neurology_sales",  143.8),
                        data.get("_neurology_budget", 138.0)),
    ]
    total_lbe = sum(t[1] for t in tas)

    add_textbox(slide, "Net Sales by Therapeutic Area  (EUR millions)",
                0.3, 1.2, 7, 0.35, font_size=11, bold=True, color=NAVY)

    bar_colors = [BLUE, TEAL, RGBColor(0x8B, 0x5C, 0xF6)]
    max_val = 350
    bar_top  = 1.65
    bar_area_h = 4.2
    bar_w    = 1.3
    gap      = 0.55

    for i, (ta, lbe_v, bud_v) in enumerate(tas):
        xl = 0.35 + i * (bar_w * 2 + gap + 0.1)

        # Budget bar (lighter)
        bh = (bud_v / max_val) * bar_area_h
        add_rect(slide, xl, bar_top + bar_area_h - bh,
                 bar_w, bh, LGREY, MGREY, 0.5)
        add_textbox(slide, f"Bgt\n{fmt(bud_v)}",
                    xl, bar_top + bar_area_h - bh - 0.5, bar_w, 0.45,
                    font_size=8, color=DGREY, align=PP_ALIGN.CENTER)

        # LBE bar
        lh = (lbe_v / max_val) * bar_area_h
        add_rect(slide, xl + bar_w + 0.05,
                 bar_top + bar_area_h - lh, bar_w, lh, bar_colors[i])
        add_textbox(slide, f"LBE\n{fmt(lbe_v)}",
                    xl + bar_w + 0.05,
                    bar_top + bar_area_h - lh - 0.5, bar_w, 0.45,
                    font_size=8, color=DARK, align=PP_ALIGN.CENTER)

        # TA label
        add_textbox(slide, ta, xl, bar_top + bar_area_h + 0.05,
                    bar_w * 2 + 0.05, 0.35,
                    font_size=10, bold=True, color=NAVY,
                    align=PP_ALIGN.CENTER)

    # Summary table right side — positioned lower to avoid overlap with header
    cols = [
        ("LBE FY2026", "LBE FY2026",    False, False),
        ("vs Budget",   "Var vs Budget", True,  False),
        ("vs Prior LBE","Var vs Prior LBE",True, False),
        ("vs PY",       "Var vs Prior Year",True,False),
    ]
    draw_pnl_table(slide, {
        k: v for k, v in data.items()
        if k in ("Net Sales", "Distribution Margin")
    }, cols, top=2.55, left=7.85,
    col_widths=[2.1, 0.78, 0.82, 0.78, 0.78],
    title="Sales Summary")


def build_slide3_vs_prior_lbe(prs, data):
    """Slide 3: P&L vs Prior LBE."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "P&L vs Prior LBE",
                 "Movement from last month's forecast (EUR millions)")
    cols = [
        ("LBE FY2026",  "LBE FY2026",      False, False),
        ("Prior LBE",   "Prior LBE",        False, False),
        ("Variance",    "Var vs Prior LBE", True,  False),
    ]
    draw_pnl_table(slide, data, cols, top=1.25, left=0.5,
                   col_widths=[3.8, 2.6, 2.6, 3.33])


def build_slide4_vs_budget(prs, data):
    """Slide 4: P&L vs Budget."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "P&L vs Budget",
                 "FY2026 performance against approved plan (EUR millions)")
    cols = [
        ("LBE FY2026", "LBE FY2026",   False, False),
        ("Budget",     "Budget FY2026", False, False),
        ("Variance",   "Var vs Budget", True,  False),
    ]
    draw_pnl_table(slide, data, cols, top=1.25, left=0.5,
                   col_widths=[3.8, 2.6, 2.6, 3.33])


def build_slide5_vs_py(prs, data):
    """Slide 5: P&L vs Prior Year."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "P&L vs Prior Year",
                 "FY2026 LBE vs FY2025 Actuals (EUR millions)")
    cols = [
        ("LBE FY2026",  "LBE FY2026",       False, False),
        ("Prior Year",  "Prior Year (FY2025)", False, False),
        ("Variance",    "Var vs Prior Year", True,   False),
    ]
    draw_pnl_table(slide, data, cols, top=1.25, left=0.5,
                   col_widths=[3.8, 2.6, 2.6, 3.33])


def build_deck(data: dict, narrative: str, output_path: str,
               unit_data: list = None):
    """
    data: consolidated P&L dict from consolidator
    narrative: string from LLM
    unit_data: list of parsed unit dicts (with product volumes) from executor_consolidation
    """
    prs = Presentation()
    prs.slide_width  = SLD_W
    prs.slide_height = SLD_H

    # Derive per-TA sales totals from unit_data for slide 3 bar chart
    if unit_data:
        for unit in unit_data:
            key = unit["unit"].lower().replace(" ", "_")
            ns = unit["lines"].get("Net Sales", {})
            data[f"_{key}_sales"]  = ns.get("lbe", 0)
            data[f"_{key}_budget"] = ns.get("budget", 0)

    build_slide1_summary(prs, data, narrative)
    build_slide2_volumes(prs, unit_data or [])
    build_slide2_sales(prs, data)
    build_slide3_vs_prior_lbe(prs, data)
    build_slide4_vs_budget(prs, data)
    build_slide5_vs_py(prs, data)

    prs.save(output_path)
    return output_path
